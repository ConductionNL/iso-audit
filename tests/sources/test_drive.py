"""Tests voor `iso_audit.sources.drive` — DriveSource + legacy API, gws gemockt."""

from __future__ import annotations

from typing import Any
from unittest.mock import patch

import pytest

from iso_audit.sources import drive
from iso_audit.sources.base import Document


@pytest.fixture(autouse=True)
def _clean_env(monkeypatch: pytest.MonkeyPatch) -> None:
    for v in drive.FOLDER_ENV_VARS:
        monkeypatch.delenv(v, raising=False)


# ---------- _resolve_folder_id ----------


def test_resolve_expliciet() -> None:
    assert drive._resolve_folder_id("abc123") == "abc123"


def test_resolve_strip_query_params() -> None:
    assert drive._resolve_folder_id("abc123?hl=nl") == "abc123"


def test_resolve_env_eerste_variabele(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("AUDIT_SOURCE_FOLDER_ID", "uit-source")
    monkeypatch.setenv("AUDIT_DRIVE_FOLDER_ID", "uit-drive")
    # Eerste variabele in FOLDER_ENV_VARS wint.
    assert drive._resolve_folder_id() == "uit-source"


def test_resolve_fallback_naar_tweede_var(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("AUDIT_DRIVE_FOLDER_ID", "fallback")
    assert drive._resolve_folder_id() == "fallback"


def test_resolve_zonder_env_raised() -> None:
    with pytest.raises(OSError, match="Geen Drive-map"):
        drive._resolve_folder_id()


# ---------- _resolve_folder_ids (multi-folder) ----------


def test_resolve_ids_komma_sep_string() -> None:
    """Komma-gescheiden string wordt naar lijst gesplitst."""
    assert drive._resolve_folder_ids("a,b,c") == ["a", "b", "c"]


def test_resolve_ids_lijst_argument() -> None:
    """Lijst-argument wordt direct doorgegeven."""
    assert drive._resolve_folder_ids(["a", "b"]) == ["a", "b"]


def test_resolve_ids_beide_env_vars_samengevoegd(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Beide env-vars met verschillende waarden → samengevoegd, dedup, volgorde behouden."""
    monkeypatch.setenv("AUDIT_SOURCE_FOLDER_ID", "0AAP-shared")
    monkeypatch.setenv("AUDIT_DRIVE_FOLDER_ID", "1YJoG-folder")
    assert drive._resolve_folder_ids() == ["0AAP-shared", "1YJoG-folder"]


def test_resolve_ids_dedup(monkeypatch: pytest.MonkeyPatch) -> None:
    """Dezelfde ID in beide env-vars verschijnt maar één keer."""
    monkeypatch.setenv("AUDIT_SOURCE_FOLDER_ID", "samepath")
    monkeypatch.setenv("AUDIT_DRIVE_FOLDER_ID", "samepath")
    assert drive._resolve_folder_ids() == ["samepath"]


def test_resolve_ids_komma_in_env(monkeypatch: pytest.MonkeyPatch) -> None:
    """Komma-sep binnen één env-var wordt ook gesplitst."""
    monkeypatch.setenv("AUDIT_SOURCE_FOLDER_ID", "a, b ,c")
    assert drive._resolve_folder_ids() == ["a", "b", "c"]


# ---------- _is_uitgesloten ----------


@pytest.mark.parametrize(
    "naam, uitgesloten",
    [
        ("NEN-EN-ISO 9001 NL.pdf", True),
        ("ISO_IEC_27001-2022.docx", True),
        ("About the Sample Files.txt", True),
        ("Beleid Conduction.docx", False),
        ("", False),
    ],
)
def test_is_uitgesloten(naam: str, uitgesloten: bool) -> None:
    assert drive._is_uitgesloten(naam) is uitgesloten


# ---------- DriveSource init + properties ----------


def test_drivesource_init_shared_drive() -> None:
    src = drive.DriveSource(folder_id="0A1234567890")
    assert src.folder_id == "0A1234567890"
    assert src.drive_id == "0A1234567890"
    assert src.naam == "drive"


def test_drivesource_init_reguliere_map() -> None:
    src = drive.DriveSource(folder_id="1abc")
    assert src.folder_id == "1abc"
    assert src.drive_id is None


def test_drivesource_init_query_strip() -> None:
    src = drive.DriveSource(folder_id="1abc?usp=share")
    assert src.folder_id == "1abc"


def test_drivesource_init_env_fallback(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("AUDIT_SOURCE_FOLDER_ID", "env-id")
    src = drive.DriveSource()
    assert src.folder_id == "env-id"


# ---------- list_documents ----------


def test_list_documents_yields_alleen_ondersteunde() -> None:
    bestanden = [
        {
            "id": "f1",
            "name": "Beleid.docx",
            "mimeType": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
            "modifiedTime": "2026-01-01T00:00:00Z",
        },
        {
            "id": "f2",
            "name": "Doc.gdoc",
            "mimeType": "application/vnd.google-apps.document",
            "modifiedTime": "2026-02-02T00:00:00Z",
        },
        # Niet leesbaar (afbeelding, geen OCR):
        {"id": "f3", "name": "afb.png", "mimeType": "image/png", "modifiedTime": ""},
        # PDF wordt sinds 2026-08-18 wél gelezen — 91 van de 213 ongelezen bestanden.
        {"id": "f4", "name": "rapport.pdf", "mimeType": "application/pdf", "modifiedTime": ""},
        # Uitgesloten op naam:
        {
            "id": "f5",
            "name": "ISO_IEC_27001.docx",
            "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "modifiedTime": "",
        },
    ]
    src = drive.DriveSource(folder_id="x")
    with patch.object(drive, "drive_lijst_bestanden", return_value=bestanden):
        docs = list(src.list_documents())
    ids = {d.id for d in docs}
    assert ids == {"f1", "f2", "f4"}


def test_list_documents_geeft_metadata_door() -> None:
    bestanden = [
        {
            "id": "f1",
            "name": "Beleid.docx",
            "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "modifiedTime": "2026-01-01T00:00:00Z",
        }
    ]
    src = drive.DriveSource(folder_id="x")
    with patch.object(drive, "drive_lijst_bestanden", return_value=bestanden):
        doc = next(iter(src.list_documents()))
    assert doc.id == "f1"
    assert doc.titel == "Beleid.docx"
    assert doc.bron == "drive"
    assert doc.type == "docx"
    assert doc.laatst_gewijzigd == "2026-01-01T00:00:00Z"
    assert doc.inhoud_uri == "f1"


def test_list_documents_lege_modifiedtime() -> None:
    """`modifiedTime` ontbrekend → lege string in `laatst_gewijzigd`."""
    bestanden = [
        {"id": "f1", "name": "x.txt", "mimeType": "text/plain"},
    ]
    src = drive.DriveSource(folder_id="y")
    with patch.object(drive, "drive_lijst_bestanden", return_value=bestanden):
        doc = next(iter(src.list_documents()))
    assert doc.laatst_gewijzigd == ""


# ---------- multi-folder ----------


def test_list_documents_multi_folder_unie() -> None:
    """Twee folders met disjoint files → union van docs."""
    folder_a = [
        {
            "id": "a1",
            "name": "A.docx",
            "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        },
    ]
    folder_b = [
        {
            "id": "b1",
            "name": "B.docx",
            "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        },
    ]
    src = drive.DriveSource(folder_id=["foldA", "foldB"])

    def _fake_list(fid: str, drive_id: str | None = None) -> list[dict[str, Any]]:
        return folder_a if fid == "foldA" else folder_b

    with patch.object(drive, "drive_lijst_bestanden", side_effect=_fake_list):
        ids = {d.id for d in src.list_documents()}
    assert ids == {"a1", "b1"}


def test_list_documents_multi_folder_dedup_op_file_id() -> None:
    """Hetzelfde file-id in beide folders → maar één keer in output."""
    overlap = [
        {
            "id": "same-1",
            "name": "Doc.docx",
            "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        },
    ]
    src = drive.DriveSource(folder_id=["foldA", "foldB"])
    with patch.object(drive, "drive_lijst_bestanden", return_value=overlap):
        docs = list(src.list_documents())
    assert len(docs) == 1
    assert docs[0].id == "same-1"


def test_drivesource_shared_drive_detection_per_folder() -> None:
    """`0A`-prefix wordt per folder als shared-drive-id behandeld."""
    src = drive.DriveSource(folder_id=["0A-shared", "1-regular"])
    # Property `folder_ids` levert beide; `folder_id`/`drive_id` retro-compat naar de eerste.
    assert src.folder_ids == ["0A-shared", "1-regular"]
    assert src.folder_id == "0A-shared"
    assert src.drive_id == "0A-shared"  # eerste is shared
    # Verifieer dat de tweede regular als folder-only is geregistreerd.
    assert src._drive_id_voor["1-regular"] is None


# ---------- fetch_content ----------


def test_fetch_content_google_doc() -> None:
    src = drive.DriveSource(folder_id="x")
    doc = Document(
        id="d1",
        titel="Doc",
        bron="drive",
        type="google_doc",
        laatst_gewijzigd="",
        inhoud_uri="d1",
    )
    with patch.object(drive, "drive_exporteer_google_doc", return_value="text!") as mock:
        out = src.fetch_content(doc)
    assert out == "text!"
    mock.assert_called_once_with("d1")


def test_fetch_content_plain_text() -> None:
    src = drive.DriveSource(folder_id="x")
    doc = Document(
        id="d1",
        titel="x.txt",
        bron="drive",
        type="txt",
        laatst_gewijzigd="",
        inhoud_uri="d1",
    )
    with patch.object(drive, "drive_download_bestand", return_value=b"hello"):
        out = src.fetch_content(doc)
    assert out == "hello"


def test_fetch_content_andere_bron_raised() -> None:
    src = drive.DriveSource(folder_id="x")
    doc = Document(
        id="d1",
        titel="x",
        bron="jira",
        type="google_doc",
        laatst_gewijzigd="",
        inhoud_uri="d1",
    )
    with pytest.raises(ValueError, match="DriveSource"):
        src.fetch_content(doc)


def test_fetch_content_onbekend_type_raised() -> None:
    src = drive.DriveSource(folder_id="x")
    doc = Document(
        id="d1",
        titel="x",
        bron="drive",
        type="onbekend",
        laatst_gewijzigd="",
        inhoud_uri="d1",
    )
    with pytest.raises(ValueError, match="Onbekend Document-type"):
        src.fetch_content(doc)


# ---------- list_findings + healthcheck ----------


def test_list_findings_geeft_lege_iterator() -> None:
    src = drive.DriveSource(folder_id="x")
    assert list(src.list_findings("sessie-1")) == []


def test_healthcheck_ok() -> None:
    src = drive.DriveSource(folder_id="x")
    with patch.object(drive, "drive_lijst_bestanden", return_value=[{"id": "a"}]):
        h = src.healthcheck()
    assert h["status"] == "ok"
    assert h["naam"] == "drive"
    assert h["aantal_bestanden"] == 1


def test_healthcheck_fail_geeft_een_soort_en_lekt_de_ruwe_melding_niet() -> None:
    """De melding van de leverancier hoort in het serverlog, niet in het antwoord.

    Deze test asserteerde eerder dat "boom" in `reden` stond — dat legde het lek vast als
    gewenst gedrag. De ruwe tekst kan een URL met credential of een responsbody bevatten.
    """
    src = drive.DriveSource(folder_id="x")
    with patch.object(drive, "drive_lijst_bestanden", side_effect=RuntimeError("boom")):
        h = src.healthcheck()
    assert h["status"] == "fail"
    assert h["soort"]
    assert "boom" not in str(h["reden"])


def test_healthcheck_multi_folder_aggregeert() -> None:
    """Healthcheck telt bestanden per folder + totaal."""
    src = drive.DriveSource(folder_id=["foldA", "foldB"])

    def _fake_list(fid: str, drive_id: str | None = None) -> list[dict[str, Any]]:
        return [{"id": f"{fid}-1"}, {"id": f"{fid}-2"}] if fid == "foldA" else [{"id": "b-1"}]

    with patch.object(drive, "drive_lijst_bestanden", side_effect=_fake_list):
        h = src.healthcheck()
    assert h["status"] == "ok"
    assert h["aantal_bestanden"] == 3
    assert h["per_folder"] == {"foldA": 2, "foldB": 1}


def test_healthcheck_multi_folder_fail_eerste_folder() -> None:
    """Eerste falende folder → status=fail met de specifieke folder benoemd."""
    src = drive.DriveSource(folder_id=["foldA", "foldB"])

    def _fake_list(fid: str, drive_id: str | None = None) -> list[dict[str, Any]]:
        if fid == "foldB":
            raise RuntimeError("permission denied op foldB")
        return [{"id": "a-1"}]

    with patch.object(drive, "drive_lijst_bestanden", side_effect=_fake_list):
        h = src.healthcheck()
    assert h["status"] == "fail"
    # De folder staat in `tenant` (die hoort bewust in de trail), niet in de vrije tekst:
    # daar zou de ruwe leveranciersmelding mee naar binnen komen.
    assert h["tenant"] == "foldB"
    assert "permission denied" not in str(h["reden"])


# ---------- Registry-registratie ----------


def test_drivesource_geregistreerd() -> None:
    """DriveSource zou via @register beschikbaar moeten zijn in SourceRegistry."""
    from iso_audit.sources import available, get

    assert "drive" in available()
    assert get("drive") is drive.DriveSource


# ---------- Legacy haal_documenten_op ----------


def test_haal_documenten_op_lege_lijst_raised() -> None:
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=[]),
        pytest.raises(RuntimeError, match="Geen bestanden"),
    ):
        drive.haal_documenten_op(folder_id="x")


def test_haal_documenten_op_happy_path() -> None:
    bestanden = [
        {"id": "f1", "name": "x.txt", "mimeType": "text/plain", "modifiedTime": "2026-01-01"},
    ]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=b"content"),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert len(docs) == 1
    assert docs[0]["tekst"] == "content"
    assert review == []


def test_haal_documenten_op_niet_tekstueel_naar_review() -> None:
    bestanden = [
        {"id": "f1", "name": "afb.png", "mimeType": "image/png", "modifiedTime": ""},
    ]
    with patch.object(drive, "drive_lijst_bestanden", return_value=bestanden):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert docs == []
    assert len(review) == 1
    assert "image/png" in review[0]["reden"]


def test_haal_documenten_op_leesfout_naar_review() -> None:
    bestanden = [
        {
            "id": "f1",
            "name": "Doc.gdoc",
            "mimeType": "application/vnd.google-apps.document",
            "modifiedTime": "",
        },
    ]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_exporteer_google_doc", side_effect=RuntimeError("oops")),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert docs == []
    assert review[0]["reden"].startswith("Leesfout")


# ---------- probe: status per locatie ----------
#
# Tot 2026-08-17 slaagde `probe()` zodra de API-aanroep lukte. De Drive-query is
# `'<id>' in parents`, dus een bestand-ID matcht niets — geen fout, een lege lijst — en de
# UI meldde **gekoppeld** terwijl elke run nul documenten uit die locatie las.


def _probe_met(
    ids: str,
    *,
    info: dict[str, dict[str, str] | None],
    telling: dict[str, tuple[int, bool]],
) -> dict[str, Any]:
    """Draai `probe()` met de client gestubd per locatie-ID."""
    src = drive.DriveSource(folder_id=ids)
    with (
        patch.object(drive, "drive_locatie_info", side_effect=lambda fid: info.get(fid)),
        patch.object(
            drive,
            "drive_inhoud_telling",
            side_effect=lambda fid, drive_id=None: telling[fid],
        ),
    ):
        return src.probe()


_MAP = "application/vnd.google-apps.folder"
_PDF = "application/pdf"


def test_probe_bestand_id_is_geen_koppeling() -> None:
    """Een bestand-ID levert nul documenten op; dat mag geen groen bolletje geven."""
    uit = _probe_met(
        "bestand1",
        info={"bestand1": {"id": "bestand1", "naam": "Beleid.pdf", "mime": _PDF}},
        telling={"bestand1": (0, False)},
    )
    assert uit["status"] == "fail"
    rij = uit["locaties"][0]
    assert rij["soort"] == "geen-map"
    assert rij["status"] == "leeg"
    assert "geen map" in str(rij["reden"])


def test_probe_lege_map_beweert_geen_oorzaak() -> None:
    """Bij een echt lege map is de oorzaak niet vast te stellen — dan niets beweren."""
    uit = _probe_met(
        "map1",
        info={"map1": {"id": "map1", "naam": "Interne audits", "mime": _MAP}},
        telling={"map1": (0, False)},
    )
    rij = uit["locaties"][0]
    assert rij["soort"] == "map"
    assert rij["status"] == "leeg"
    assert "geen map" not in str(rij["reden"])
    assert "geen bestanden of submappen" in str(rij["reden"])


def test_probe_map_met_alleen_submappen_is_niet_leeg() -> None:
    """Nul bestanden maar wél submappen: een recursieve run leest daar wél uit."""
    uit = _probe_met(
        "map1",
        info={"map1": {"id": "map1", "naam": "Beleid", "mime": _MAP}},
        telling={"map1": (0, True)},
    )
    assert uit["status"] == "ok"
    assert uit["locaties"][0]["status"] == "ok"


def test_probe_een_goede_en_een_lege_locatie_blijft_gekoppeld() -> None:
    """Eén verkeerd geplakt ID mag een werkende configuratie niet rood maken."""
    uit = _probe_met(
        "0AAP-shared,bestand1",
        info={
            "0AAP-shared": {"id": "0AAP-shared", "naam": "Conduction", "mime": ""},
            "bestand1": {"id": "bestand1", "naam": "Beleid.pdf", "mime": _PDF},
        },
        telling={"0AAP-shared": (409, True), "bestand1": (0, False)},
    )
    assert uit["status"] == "ok"
    per_status = {str(r["id"]): r["status"] for r in uit["locaties"]}
    assert per_status == {"0AAP-shared": "ok", "bestand1": "leeg"}


def test_probe_shared_drive_wordt_als_zodanig_herkend() -> None:
    """`files.get` geeft voor een Shared Drive-root geen map-mime; het ID-prefix wel."""
    uit = _probe_met(
        "0AAP-shared",
        info={"0AAP-shared": {"id": "0AAP-shared", "naam": "Conduction", "mime": ""}},
        telling={"0AAP-shared": (12, False)},
    )
    rij = uit["locaties"][0]
    assert rij["soort"] == "shared-drive"
    assert rij["naam"] == "Conduction"


def test_probe_zonder_naam_blijft_bruikbaar() -> None:
    """De naam is comfort, geen voorwaarde: zonder naam blijft de locatie gekoppeld."""
    uit = _probe_met(
        "map1",
        info={"map1": None},
        telling={"map1": (3, False)},
    )
    assert uit["status"] == "ok"
    rij = uit["locaties"][0]
    assert rij["naam"] == ""
    assert rij["soort"] == "onbekend"
    assert rij["status"] == "ok"


def test_probe_verbindingsfout_weegt_zwaarder_dan_leeg() -> None:
    """Een geweigerde credential zegt iets anders dan een lege map; die moet vooropstaan."""
    src = drive.DriveSource(folder_id="map1,map2")

    def _telling(fid: str, drive_id: str | None = None) -> tuple[int, bool]:
        if fid == "map2":
            raise RuntimeError("403 forbidden")
        return (0, False)

    with (
        patch.object(drive, "drive_locatie_info", return_value=None),
        patch.object(drive, "drive_inhoud_telling", side_effect=_telling),
    ):
        uit = src.probe()

    assert uit["status"] == "fail"
    assert uit["tenant"] == "map2"
    assert uit["soort"] == "auth"


def test_probe_telt_niet_recursief() -> None:
    """De statusregel mag geen enumeratie worden — één telling per locatie, niet meer."""
    src = drive.DriveSource(folder_id="map1,map2")
    aanroepen: list[str] = []

    def _telling(fid: str, drive_id: str | None = None) -> tuple[int, bool]:
        aanroepen.append(fid)
        return (5, True)

    with (
        patch.object(drive, "drive_locatie_info", return_value=None),
        patch.object(drive, "drive_inhoud_telling", side_effect=_telling),
        patch.object(drive, "drive_lijst_bestanden", side_effect=AssertionError("recursief!")),
    ):
        src.probe()

    assert aanroepen == ["map1", "map2"]


# Type-cast voor mypy zodat tests/typing klopt.
_ = Any


# ---------- dekking: wat niet gelezen wordt, wordt gemeld ----------
#
# Tot 2026-08-18 verdwenen 92 van de 512 bestanden in de gekoppelde Shared Drive via
# `logger.debug("Skip (onbekend MIME)")`: niet in het gemelde aantal voor handmatige review,
# en op INFO-niveau geen enkele regel. Deze tests zijn de gate daarop.


def _xlsx_bytes(bladen: dict[str, list[list[Any]]]) -> bytes:
    """Bouw een xlsx in het geheugen. Geen binaire fixture in de repo: wat erin zit moet
    leesbaar zijn in de test die het gebruikt."""
    import io

    import openpyxl

    boek = openpyxl.Workbook()
    boek.remove(boek.active)
    for naam, rijen in bladen.items():
        blad = boek.create_sheet(naam)
        for rij in rijen:
            blad.append(rij)
    buffer = io.BytesIO()
    boek.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(regels: list[str]) -> bytes:
    import io

    import pptx

    presentatie = pptx.Presentation()
    dia = presentatie.slides.add_slide(presentatie.slide_layouts[5])
    dia.shapes.title.text = regels[0]
    buffer = io.BytesIO()
    presentatie.save(buffer)
    return buffer.getvalue()


def _pdf_bytes(tekst: str | None) -> bytes:
    """Minimale, met de hand opgebouwde PDF — met of zonder tekstlaag.

    Met de hand en niet als gecommitte binary: een auditor moet in de test kunnen zien wat
    erin zit. `tekst=None` levert een pagina zonder tekstlaag op, wat een gescande PDF
    nabootst: `extract_text()` geeft dan nul tekens.
    """
    stroom = b"" if tekst is None else f"BT /F1 12 Tf 72 720 Td ({tekst}) Tj ET".encode("latin-1")
    objecten = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length %d >>\nstream\n" % len(stroom) + stroom + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    uit = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for nummer, obj in enumerate(objecten, start=1):
        offsets.append(len(uit))
        uit += b"%d 0 obj\n" % nummer + obj + b"\nendobj\n"
    xref = len(uit)
    uit += b"xref\n0 %d\n" % (len(objecten) + 1)
    uit += b"0000000000 65535 f \n"
    for offset in offsets:
        uit += b"%010d 00000 n \n" % offset
    uit += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objecten) + 1,
        xref,
    )
    return bytes(uit)


def _bestand(id_: str, naam: str, mime: str) -> dict[str, Any]:
    return {"id": id_, "name": naam, "mimeType": mime, "modifiedTime": "2026-05-05T00:00:00Z"}


def test_markdown_wordt_gelezen_als_tekst() -> None:
    """`Auditrapport_beide_v3.3_2026-05-05.md` bleef buiten het landschap omdat markdown een
    andere tekst-MIME heeft dan `text/plain` — niet omdat het onleesbaar was."""
    bestanden = [_bestand("f1", "Auditrapport.md", "text/markdown")]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=b"# Auditrapport"),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert [d["tekst"] for d in docs] == ["# Auditrapport"]
    assert review == []


def test_html_en_csv_worden_gelezen() -> None:
    bestanden = [
        _bestand("f1", "index.html", "text/html"),
        _bestand("f2", "lijst.csv", "text/csv"),
    ]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=b"inhoud"),
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    assert len(docs) == 2


def test_xlsx_levert_celtekst_per_blad() -> None:
    """Een CSV-export zou alleen het eerste blad geven; dat is de stille onvolledigheid."""
    inhoud = _xlsx_bytes({"Acties": [["Risico", "Hoog"]], "Beoordeling": [["Q1", "akkoord"]]})
    bestanden = [_bestand("f1", "RIE.xlsx", drive.XLSX_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=inhoud),
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    tekst = docs[0]["tekst"]
    assert "Acties" in tekst and "Hoog" in tekst
    assert "Beoordeling" in tekst and "akkoord" in tekst, "het tweede blad mag niet wegvallen"


def test_pptx_levert_diatekst() -> None:
    inhoud = _pptx_bytes(["Managementreview 2026"])
    bestanden = [_bestand("f1", "Review.pptx", drive.PPTX_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=inhoud),
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    assert "Managementreview 2026" in docs[0]["tekst"]


def test_pdf_wordt_gelezen() -> None:
    """91 PDF's, waaronder de auditrapporten van de certificerende instantie en de VvT."""
    bestanden = [_bestand("f1", "Auditrapport ISMS.pdf", "application/pdf")]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=_pdf_bytes("Auditrapport ISMS")),
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    assert "Auditrapport ISMS" in docs[0]["tekst"]


def test_google_sheet_wordt_als_xlsx_geexporteerd() -> None:
    """Als CSV komt alleen het eerste blad mee; als xlsx alle bladen."""
    inhoud = _xlsx_bytes({"Blad1": [["a"]], "Blad2": [["b"]]})
    bestanden = [_bestand("f1", "Actielijst", drive.GOOGLE_SHEET_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_exporteer_bytes", return_value=inhoud) as export,
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    export.assert_called_once_with("f1", drive.XLSX_MIME)
    assert "Blad2" in docs[0]["tekst"]


def test_google_slides_wordt_als_tekst_geexporteerd() -> None:
    bestanden = [_bestand("f1", "Presentatie", drive.GOOGLE_SLIDES_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_exporteer_bytes", return_value=b"dia-tekst") as export,
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    export.assert_called_once_with("f1", "text/plain")
    assert docs[0]["tekst"] == "dia-tekst"


def test_onbekend_type_wordt_gemeld_en_niet_stil_overgeslagen(
    caplog: pytest.LogCaptureFixture,
) -> None:
    """Dit is het stille gat: geen melding, en niet in het aantal voor handmatige review."""
    bestanden = [
        _bestand("f1", "notitie.txt", "text/plain"),
        _bestand("f2", "raadsel.bin", "application/x-onbekend"),
    ]
    with (
        caplog.at_level("INFO", logger="iso_audit.sources.drive"),
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=b"inhoud"),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert len(docs) == 1
    assert [r["naam"] for r in review] == ["raadsel.bin"]
    assert "application/x-onbekend" in review[0]["reden"]
    meldingen = "\n".join(caplog.messages)
    assert "application/x-onbekend" in meldingen
    assert "Niet gelezen (1): onbekend type: application/x-onbekend" in meldingen


def test_niet_leesbaar_meldt_de_reden(caplog: pytest.LogCaptureFixture) -> None:
    """ "Onleesbaar" zonder reden stuurt de auditor net zo hard het verkeerde bos in."""
    bestanden = [_bestand("f1", "scan.png", "image/png")]
    with (
        caplog.at_level("INFO", logger="iso_audit.sources.drive"),
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
    ):
        _, review = drive.haal_documenten_op(folder_id="x")
    assert "geen OCR" in review[0]["reden"]
    assert "image/png" in "\n".join(caplog.messages)


def test_leeg_extractieresultaat_is_geen_leeg_document() -> None:
    """Een gescande PDF levert nul tekens op. Als leeg document opgenomen classificeert de
    pipeline hem als "geen bewijs" — een oordeel over iets wat niemand heeft gelezen."""
    bestanden = [_bestand("f1", "Gescand certificaat.pdf", "application/pdf")]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=_pdf_bytes(None)),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert docs == [], "geen document met lege inhoud in het landschap"
    assert "Onleesbaar" in review[0]["reden"]
    assert "scan" in review[0]["reden"]


def test_dekking_telt_gezien_gelezen_en_redenen() -> None:
    """De dekking gaat naar het run-record; het log overleeft geen podherstart."""
    bestanden = [
        _bestand("f1", "notitie.txt", "text/plain"),
        _bestand("f2", "scan.png", "image/png"),
        _bestand("f3", "film.mp4", "video/mp4"),
        _bestand("f4", "raadsel.bin", "application/x-onbekend"),
        _bestand("f5", "ISO_IEC_27001.docx", drive.DOCX_MIME),
    ]
    gezien: list[drive.Dekkingteller] = []
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=b"inhoud"),
    ):
        drive.haal_documenten_op(folder_id="x", op_dekking=gezien.append)
    teller = gezien[0]
    assert (teller.gezien, teller.gelezen, teller.niet_gelezen) == (5, 1, 4)
    assert teller.overgeslagen["onbekend type: application/x-onbekend"] == 1
    assert teller.overgeslagen["referentiedocument, bewust uitgesloten"] == 1
    assert sum(teller.overgeslagen.values()) == teller.niet_gelezen


def test_list_documents_meldt_de_dekking(caplog: pytest.LogCaptureFixture) -> None:
    """Ook het Source-protocol-pad mag niets stil laten vallen."""
    bestanden = [
        _bestand("f1", "notitie.txt", "text/plain"),
        _bestand("f2", "raadsel.bin", "application/x-onbekend"),
    ]
    src = drive.DriveSource(folder_id="x")
    with (
        caplog.at_level("INFO", logger="iso_audit.sources.drive"),
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
    ):
        docs = list(src.list_documents())
    assert len(docs) == 1
    meldingen = "\n".join(caplog.messages)
    assert "2 bestand(en) gezien, 1 gelezen, 1 niet gelezen" in meldingen
    assert "onbekend type: application/x-onbekend" in meldingen


def test_doel_van_snelkoppeling_telt_een_keer() -> None:
    """De client lost een snelkoppeling op naar het doel-record. Zit dat doel óók rechtstreeks
    in scope, dan komt het twee keer in de lijst — de dedup op file-id vangt dat."""
    doel = _bestand("echt", "VvT Conduction ISO 27001.pdf", "application/pdf")
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=[doel, dict(doel)]),
        patch.object(drive, "drive_download_bestand", return_value=_pdf_bytes("VvT")),
    ):
        docs, _ = drive.haal_documenten_op(folder_id="x")
    assert [d["id"] for d in docs] == ["echt"]


def test_docx_tabellen_worden_gelezen() -> None:
    """Een actiepuntenlijst is één tabel en levert nul alinea's op.

    Gemeten in de eerste productierun (2026-08-21): `Actiepunten uit Waveland.docx` kwam
    binnen als "geen tekst, mogelijk een scan" — wat voor een docx onmogelijk is.
    """
    import io as _io

    import docx as _docx

    d = _docx.Document()
    tabel = d.add_table(rows=2, cols=2)
    tabel.cell(0, 0).text = "Actiepunt"
    tabel.cell(1, 0).text = "RI&E afronden"
    buffer = _io.BytesIO()
    d.save(buffer)

    bestanden = [_bestand("f1", "Actiepunten.docx", drive.DOCX_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=buffer.getvalue()),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert docs, "een docx met alleen een tabel mag niet als leeg gelden"
    assert "RI&E afronden" in docs[0]["tekst"]
    assert review == []


def test_reden_bij_geen_tekst_hangt_af_van_het_formaat() -> None:
    """ "Mogelijk een scan" hoort bij PDF. Bij een docx wijst dat de auditor de verkeerde
    kant op — juist bij het bestand dat wél te repareren is."""
    import io as _io

    import docx as _docx

    leeg = _io.BytesIO()
    _docx.Document().save(leeg)

    bestanden = [
        _bestand("f1", "Leeg.docx", drive.DOCX_MIME),
        _bestand("f2", "Scan.pdf", "application/pdf"),
    ]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(
            drive, "drive_download_bestand", side_effect=[leeg.getvalue(), _pdf_bytes(None)]
        ),
    ):
        _, review = drive.haal_documenten_op(folder_id="x")
    per_naam = {r["naam"]: r["reden"] for r in review}
    assert "scan" not in per_naam["Leeg.docx"]
    assert "tekstvakken" in per_naam["Leeg.docx"]
    assert "scan" in per_naam["Scan.pdf"]


def test_docx_met_alleen_afbeeldingen_meldt_dat_als_feit() -> None:
    """Gemeten op 2026-08-21: `Actiepunten uit Waveland.docx` is 569 KB met drie lege
    alinea's, nul tabellen en zes `w:drawing`-elementen — screenshots in een Word-bestand.
    "Mogelijk een scan" zei daar het verkeerde over."""
    import io as _io

    import docx as _docx

    # Een `w:drawing` rechtstreeks in de body: `add_picture` vraagt een echte
    # afbeeldingsheader, en die hoort niet in een test over tekstextractie.
    d = _docx.Document()
    alinea = d.add_paragraph()
    alinea._p.append(
        _docx.oxml.parse_xml(
            '<w:drawing xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
        )
    )
    buffer = _io.BytesIO()
    d.save(buffer)

    bestanden = [_bestand("f1", "Actiepunten.docx", drive.DOCX_MIME)]
    with (
        patch.object(drive, "drive_lijst_bestanden", return_value=bestanden),
        patch.object(drive, "drive_download_bestand", return_value=buffer.getvalue()),
    ):
        docs, review = drive.haal_documenten_op(folder_id="x")
    assert docs == []
    assert "ingevoegde afbeelding" in review[0]["reden"]
    assert "scan" not in review[0]["reden"]
    assert "OCR" in review[0]["reden"]
