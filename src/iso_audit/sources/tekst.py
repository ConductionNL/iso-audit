"""Tekst uit bestandsbytes — gedeeld door elke bron die bestanden leest.

## Waarom dit een eigen module is

Deze lezers stonden in `sources/drive.py` en raken Drive niet: het zijn functies van bytes naar
tekst. Toen `nextcloud-bron` erbij kwam, was de keuze een tweede set lezers of één gedeelde.

Een tweede set maakt van "leest het tool xlsx-tabellen?" een vraag met twee antwoorden, die na
één wijziging uit elkaar lopen. Dat is precies het soort stille verschil waar de
dekkingsmeldingen van 2026-08-18 voor bestaan — dus één set, hier.

Wat **niet** hier hoort: de Google-exports (Docs, Sheets, Slides). Die vragen een Drive-API en
zijn dus wél bronspecifiek.

## De regel die deze module draagt

Een geslaagde extractie die nul tekens oplevert, is een **storing** en geen leeg document. Een
gescande PDF als leeg document opgenomen classificeert de pipeline als "geen bewijs" — een
oordeel over iets wat niemand heeft gelezen, op een clausule waar het bewijs bestaat.
"""

from __future__ import annotations

import io
import re
import xml.etree.ElementTree as ET  # nosec B405 — zie OnleesbaarDocumentError
import zipfile
from collections.abc import Callable

import docx
import openpyxl
import pptx
import pypdf

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

ODF_MIMES: dict[str, str] = {
    "application/vnd.oasis.opendocument.text": "odt",
    "application/vnd.oasis.opendocument.spreadsheet": "ods",
    "application/vnd.oasis.opendocument.presentation": "odp",
    "application/vnd.oasis.opendocument.graphics": "odg",
}
"""De vier OpenDocument-formaten die op een gedeelde schijf voorkomen.

Alle vier gaan door dezelfde lezer: het verschil tussen odt en odp zit in de omhulling
(`office:text` tegen `draw:page`), niet in waar de tekst staat. Een aparte lezer per formaat
zou drie keer dezelfde wandeling zijn met drie keer een eigen kans op afwijken.

`odg` is een tekening. Die levert vaak niets op, en dat is een uitkomst en geen fout — de
caller meldt hem dan als onleesbaar, net als een gescande pdf."""

GESCANDE_FORMATEN: frozenset[str] = frozenset({"application/pdf"})
"""Formaten waarbij "geen tekst" op een scan kan duiden.

Voor de rest kan het dat niet. Een `.docx` zonder tekst is geen scan maar een document
waarvan de inhoud buiten het bereik van de lezer valt — tot 2026-08-21 stond in de melding
"mogelijk een scan of een leeg bestand" voor élk formaat, en dat wees de auditor de verkeerde
kant op bij precies het bestand dat wél te repareren was."""


MAX_ODF_INHOUD = 32 * 1024 * 1024
"""Bovengrens op de **uitgepakte** `content.xml`.

Op de zip meten heeft geen zin: die kant is juist het probleem. 40 MB nullen comprimeert tot
enkele kilobytes, dus een grens op het bestand laat een zip-bom door. Zelfde grootte als
`MAX_ANTWOORD` in de WebDAV-client, om dezelfde reden: één getal om te verantwoorden."""

_ODF_DOCTYPE = re.compile(r"<!DOCTYPE", re.IGNORECASE)
_DOCTYPE_VENSTER = 2048
"""Zie `clients/nextcloud.py`: een DTD staat in de prolog, verderop is dit gewoon tekst."""


class OnleesbaarDocumentError(Exception):
    """Het bestand kon niet worden uitgelezen — kapot, of geweigerd om veiligheidsredenen.

    Bewust onderscheiden van `LeegDocumentError`: leeg betekent "gelezen, niets gevonden" en
    onleesbaar betekent "niet gelezen". Dat verschil staat in de dekkingsmelding aan de auditor,
    en één fout voor beide zou die melding onbruikbaar maken.
    """


class LeegDocumentError(Exception):
    """De extractie lukte, maar leverde geen tekst op.

    Bewust een eigen fout en geen lege string: een gescande PDF levert nul tekens, en als
    document met lege inhoud opgenomen classificeert de pipeline hem als "geen bewijs" — een
    oordeel over iets wat niemand heeft gelezen, op een clausule waar het bewijs bestaat.
    Dezelfde regel als bij de classificatie sinds 2026-08-17, waar een afgekapt antwoord ook
    geen leeg oordeel meer is.
    """


def tekst_uit_docx(inhoud: bytes) -> str:
    """Alinea's **en** tabelcellen.

    De tabellen ontbraken tot 2026-08-21, en dat is geen detail: een actiepuntenlijst of een
    RI&E-overzicht is één tabel en levert dan nul alinea's op. Gemeten in de eerste
    productierun — `Actiepunten uit Waveland.docx` kwam binnen als "geen tekst, mogelijk een
    scan", wat voor een docx onmogelijk is.

    Tekstvakken en koppen/voetteksten blijven buiten bereik: die zitten niet in het
    `document.body`-model van python-docx. Een docx die daardoor leeg blijft, wordt gemeld als
    onleesbaar — zichtbaar, niet stil.
    """
    doc = docx.Document(io.BytesIO(inhoud))
    delen = [p.text for p in doc.paragraphs]
    for tabel in doc.tables:
        for rij in tabel.rows:
            cellen = [c.text.strip() for c in rij.cells if c.text.strip()]
            if cellen:
                delen.append("\t".join(cellen))
    tekst = "\n".join(delen)
    if not tekst.strip():
        # Geen tekst én tekeningen in de body: dan is de inhoud ingevoegde afbeeldingen, en
        # dat is een feit in plaats van een vermoeden. Gemeten op 2026-08-21:
        # `Actiepunten uit Waveland.docx` is 569 KB met drie lege alinea's, nul tabellen en
        # zes `w:drawing`-elementen — screenshots in een Word-bestand. "Mogelijk een scan"
        # zei daar het verkeerde over; alleen OCR zou hier iets opleveren.
        tekeningen = doc.element.body.xml.count("w:drawing")
        if tekeningen:
            raise LeegDocumentError(
                f"bevat {tekeningen} ingevoegde afbeelding(en) en geen tekst; "
                "zonder OCR is hier niets uit te lezen"
            )
    return tekst


def tekst_uit_xlsx(inhoud: bytes) -> str:
    """Celtekst per blad, met de bladnaam als kop.

    `data_only=True` geeft de laatst berekende waarde in plaats van de formule: in een
    RI&E-actielijst is "hoog" het bewijs, niet `=ALS(...)`.
    """
    boek = openpyxl.load_workbook(io.BytesIO(inhoud), data_only=True, read_only=True)
    regels: list[str] = []
    for blad in boek.worksheets:
        regels.append(f"## {blad.title}")
        for rij in blad.iter_rows(values_only=True):
            cellen = [str(c) for c in rij if c is not None and str(c).strip()]
            if cellen:
                regels.append("\t".join(cellen))
    boek.close()
    return "\n".join(regels)


def tekst_uit_pptx(inhoud: bytes) -> str:
    """Tekst per dia, inclusief tabellen — om dezelfde reden als bij docx."""
    presentatie = pptx.Presentation(io.BytesIO(inhoud))
    regels: list[str] = []
    for nummer, dia in enumerate(presentatie.slides, start=1):
        regels.append(f"## Dia {nummer}")
        for vorm in dia.shapes:
            tekst = getattr(vorm, "text", "")
            if tekst and tekst.strip():
                regels.append(tekst)
            if getattr(vorm, "has_table", False):
                for rij in vorm.table.rows:
                    cellen = [c.text.strip() for c in rij.cells if c.text.strip()]
                    if cellen:
                        regels.append("\t".join(cellen))
    return "\n".join(regels)


def tekst_uit_pdf(inhoud: bytes) -> str:
    """Doorlopende tekst per pagina via `pypdf`. Geen OCR.

    Een gescande PDF levert hier nul tekens op; die wordt door de caller als onleesbaar
    gemeld en niet als leeg document opgenomen.
    """
    lezer = pypdf.PdfReader(io.BytesIO(inhoud))
    return "\n".join(pagina.extract_text() or "" for pagina in lezer.pages)


def _odf_regels(element: ET.Element) -> list[str]:
    """Wandel de boom in documentorde en maak per alinea, kop of tabelrij één regel.

    De namespace wordt afgekapt (`{...}p` → `p`) in plaats van uitgeschreven: ODF 1.2 en 1.3
    gebruiken dezelfde lokale namen onder verschillende namespace-URI's, en een bestand van een
    oudere LibreOffice zou dan stil nul regels opleveren.
    """
    regels: list[str] = []
    for kind in element:
        naam = kind.tag.rsplit("}", 1)[-1]
        if naam in ("p", "h"):
            regel = "".join(kind.itertext()).strip()
            if regel:
                regels.append(regel)
        elif naam == "table-row":
            cellen = [c.strip() for c in ("".join(k.itertext()) for k in kind) if c.strip()]
            if cellen:
                regels.append("\t".join(cellen))
        else:
            regels.extend(_odf_regels(kind))
    return regels


def tekst_uit_odf(inhoud: bytes) -> str:
    """Tekst uit een OpenDocument-bestand (odt, ods, odp, odg) — stdlib, geen extra pakket.

    Een ODF-bestand is een zip met `content.xml`. Op de Nextcloud-canary lagen 32 van deze
    bestanden (2026-08-24), gemeld als "onbekend type": elf odt, elf odp, zes ods en vier odg.
    Dat is bijna een vijfde van die schijf, en op een schijf van een organisatie die LibreOffice
    gebruikt is het de hoofdmoot in plaats van een uitzondering.

    Geen `odfpy`: de wandeling hieronder is twintig regels, en een dependency erbij is een
    dependency om te volgen in een repo die onder 27001-scope valt.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(inhoud)) as archief:
            try:
                gegevens = archief.getinfo("content.xml")
            except KeyError as fout:
                raise OnleesbaarDocumentError(
                    "geen content.xml in het bestand; dit is geen OpenDocument-bestand"
                ) from fout
            if gegevens.file_size > MAX_ODF_INHOUD:
                raise OnleesbaarDocumentError(
                    f"content.xml van {gegevens.file_size} bytes overschrijdt de grens van "
                    f"{MAX_ODF_INHOUD}"
                )
            ruw = archief.read("content.xml")
    except zipfile.BadZipFile as fout:
        raise OnleesbaarDocumentError(f"geen leesbare zip: {fout}") from fout

    xml = ruw.decode("utf-8", errors="replace")
    if _ODF_DOCTYPE.search(xml[:_DOCTYPE_VENSTER]):
        raise OnleesbaarDocumentError(
            "content.xml bevat een DOCTYPE; geweigerd wegens entity-expansie"
        )
    try:
        wortel = ET.fromstring(xml)  # nosec B314 — DOCTYPE geweigerd en grootte begrensd, zie boven
    except ET.ParseError as fout:
        raise OnleesbaarDocumentError(f"content.xml is geen geldige XML: {fout}") from fout
    return "\n".join(_odf_regels(wortel))


_BINAIRE_LEZERS: dict[str, Callable[[bytes], str]] = {
    DOCX_MIME: tekst_uit_docx,
    XLSX_MIME: tekst_uit_xlsx,
    PPTX_MIME: tekst_uit_pptx,
    "application/pdf": tekst_uit_pdf,
    **dict.fromkeys(ODF_MIMES, tekst_uit_odf),
}
"""Per binair formaat één lezer. De rest wordt als tekst gedecodeerd."""


def lees_bytes(inhoud: bytes, mime: str) -> str:
    """Tekst uit `inhoud` volgens `mime`; onbekende types worden als tekst gedecodeerd.

    Raist geen `LeegDocumentError` — de caller beslist wat een leeg resultaat betekent, omdat
    de reden per bron verschilt: een Drive-export kan leeg zijn doordat de export faalde, een
    WebDAV-download doordat het bestand zelf leeg is.
    """
    lezer = _BINAIRE_LEZERS.get(mime)
    return lezer(inhoud) if lezer else inhoud.decode("utf-8", errors="replace")


PLATTE_TEKST: frozenset[str] = frozenset({"text/plain", "text/markdown", "text/csv", "text/html"})
"""Formaten waarbij "geen tekst" simpelweg "leeg bestand" betekent.

Er valt niets te vermoeden over tekstvakken of scans in een `.txt` van nul bytes."""


def leeg_reden(mime: str) -> str:
    """Waarom een geslaagde extractie geen tekst opleverde, per formaat.

    Eén reden voor alle formaten was juist de fout: "mogelijk een scan" bij een docx wees de
    auditor de verkeerde kant op bij precies het bestand dat wél te repareren was. En bij een
    lege `.txt` op de Nextcloud-canary (2026-08-22) stond er "mogelijk staat de inhoud in
    tekstvakken" — over een bestand van nul bytes.
    """
    if mime in GESCANDE_FORMATEN:
        return "mogelijk een scan of een leeg bestand"
    if mime in PLATTE_TEKST:
        return "het bestand is leeg"
    return "mogelijk staat de inhoud in afbeeldingen, tekstvakken of koppen"
